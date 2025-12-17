import streamlit as st
from pptx import Presentation
from openpyxl import Workbook
from openpyxl.drawing.image import Image
from openpyxl.styles import Font
import os
import tempfile

st.set_page_config(page_title="PPTX to Excel Converter", layout="centered")

st.title("📊 Convert PPTX → Excel")
st.write("Upload file PPTX, lalu download hasil Excel-nya.")

uploaded_file = st.file_uploader("Upload file PPTX", type=["pptx"])

if uploaded_file:
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, uploaded_file.name)

        with open(pptx_path, "wb") as f:
            f.write(uploaded_file.read())

        prs = Presentation(pptx_path)
        wb = Workbook()
        ws = wb.active
        ws.title = "Convert PPTX"

        row_cursor = 1
        img_counter = 1
        img_dir = os.path.join(tmpdir, "images")
        os.makedirs(img_dir, exist_ok=True)

        for slide_index, slide in enumerate(prs.slides, start=1):
            ws.merge_cells(start_row=row_cursor, start_column=1,
                           end_row=row_cursor, end_column=6)
            header = ws.cell(row=row_cursor, column=1)
            header.value = f"SLIDE {slide_index}"
            header.font = Font(bold=True)
            row_cursor += 2

            for shape in slide.shapes:

                # TEXT
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        ws.merge_cells(start_row=row_cursor, start_column=1,
                                       end_row=row_cursor, end_column=6)
                        ws.cell(row=row_cursor, column=1).value = text
                        row_cursor += 1

                # IMAGE
                if shape.shape_type == 13:  # PICTURE
                    img_path = os.path.join(img_dir, f"img_{img_counter}.png")
                    with open(img_path, "wb") as f:
                        f.write(shape.image.blob)

                    img = Image(img_path)
                    img.width = 200
                    img.height = 150
                    ws.add_image(img, f"A{row_cursor}")

                    row_cursor += 10
                    img_counter += 1

            row_cursor += 2

        output_path = os.path.join(tmpdir, "hasil_convert.xlsx")
        wb.save(output_path)

        with open(output_path, "rb") as f:
            st.success("✅ Convert berhasil!")
            st.download_button(
                label="⬇️ Download Excel",
                data=f,
                file_name="hasil_convert.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
